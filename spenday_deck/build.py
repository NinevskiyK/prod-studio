"""SpenDay — pitch deck builder. Run: python build.py -> SpenDay.pptx"""
import os
from pptx import Presentation
from deck_lib import *  # noqa: F401,F403  (design-system helpers + enums)

OUT = os.path.join(os.path.dirname(__file__), "SpenDay.pptx")


def new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def card(slide, x, y, w, h, fill=PAPER, line=LINE, lw=1.2, radius=0.18,
         shadow=True):
    return rect(slide, x, y, w, h, fill=fill, line=line, lw=lw, radius=radius,
                shadow=shadow)


def blobs_dark(slide):
    oval(slide, 9.3, -2.1, 6.6, 6.6, GLOW_D)
    oval(slide, -1.7, 4.6, 4.6, 4.6, GLOW_D)


# ============================================================== SLIDE 1 ======
def s01(prs):
    s = new(prs); base(s, INK); blobs_dark(s)
    oval(s, 8.6, 0.7, 0.5, 0.5, INK2)
    wordmark(s, MX, 0.78, size=0.42, on_dark=True)
    textbox(s, MX, 2.35, 8.0, 0.4,
            [P(R("ПРОДУКТОВАЯ СТУДИЯ · 2026", F_BOLD, 13, INDIGO_L, sp=3.0))])
    textbox(s, MX - 0.04, 2.72, 8.2, 1.5,
            [P(R("SpenDay", F_BLACK, 70, PAPER))])
    textbox(s, MX, 4.18, 8.1, 1.4,
            [P([R("Готовые сценарии досуга\n", F_XBOLD, 30, PAPER),
                R("с понятным бюджетом", F_XBOLD, 30, INDIGO_L)], ls=1.08)])
    textbox(s, MX, 5.95, 7.6, 0.6,
            [P(R("Выбрать вечер за 10 минут — а не за вечер бесконечного поиска.",
                 F_MED, 15.5, "C9C7E6"))])
    glow(s, 11.0, 4.5, 4.6, 6.2, GLOW_D)
    phone(s, "feed", height=6.5, cx=10.95, top=1.25)


# ============================================================== SLIDE 2 ======
def s02(prs):
    s = new(prs); base(s)
    header(s, "ИСТОРИЯ · ПЯТНИЦА, 19:00", "Как наша компания выбирала, куда пойти",
           2, title_size=31)
    bx, bw = MX, 6.7
    bubbles = [
        ("left",  "«Давайте в рилсы — там же куча идей»", LAV, INK),
        ("right", "«Листаем 40 минут — ничего, только залипли»", INDIGO, PAPER),
        ("left",  "«О, у нас же есть сохранённые!»", LAV, INK),
        ("right", "«Там 300 роликов про еду и ноль решений»", INDIGO, PAPER),
    ]
    y = 2.25
    for side, txt, fill, tc in bubbles:
        w = 4.85
        x = bx if side == "left" else bx + (bw - w)
        rad = 0.28
        b = rect(s, x, y, w, 0.82, fill=fill, radius=rad,
                 shadow=(side == "right"))
        set_text(b, [P(R(txt, F_MED, 14.5, tc))], anchor=MSO_ANCHOR.MIDDLE,
                 pad=0.22)
        y += 1.02
    # result panel
    px, pw = 7.95, 4.5
    card(s, px, 2.25, pw, 3.86, fill=INK, line=None, radius=0.22, shadow=True)
    textbox(s, px + 0.4, 2.62, pw - 0.8, 0.4,
            [P(R("ЧТО В ИТОГЕ", F_BOLD, 11.5, INDIGO_L, sp=2.6))])
    textbox(s, px + 0.36, 3.0, pw - 0.7, 1.2,
            [P([R("≈ 1 ", F_BLACK, 52, PAPER), R("час", F_BLACK, 30, INDIGO_L)])])
    textbox(s, px + 0.4, 4.06, pw - 0.8, 0.4,
            [P(R("обсуждений и споров в чате", F_MED, 13.5, "C9C7E6"))])
    hline(s, px + 0.4, 4.66, pw - 0.8, color=INK3, weight=1.2)
    textbox(s, px + 0.4, 4.86, pw - 0.8, 1.1,
            [P([R("Итог: ", F_BOLD, 15, PAPER),
                R("пошли туда же, куда ходим всегда.", F_MED, 15, "D8D6F0")],
               ls=1.12)])
    footer(s, 2)


# ============================================================== SLIDE 3 ======
def s03(prs):
    s = new(prs); base(s)
    header(s, "ИСТОРИЯ · ПРОБЛЕМА", "Почему текущий путь неудобный", 3,
           title_size=32)
    tiles = [
        ("Соцсети", "идея без плана", PINK),
        ("Карты", "просто поиск места", INDIGO),
        ("Отзывы", "часто устарели", AMBER),
        ("Меню", "цены — отдельно", GREEN),
        ("Сайты", "всё разрозненно", INDIGO),
        ("Чат", "вечное согласование", PINK),
    ]
    n = len(tiles); gap = 0.22
    tw = (SW - 2 * MX - gap * (n - 1)) / n
    y = 2.35
    for i, (name, cap, col) in enumerate(tiles):
        x = MX + i * (tw + gap)
        card(s, x, y, tw, 1.45, radius=0.16, shadow=True)
        oval(s, x + tw / 2 - 0.16, y + 0.26, 0.32, 0.32, col)
        textbox(s, x, y + 0.74, tw, 0.34,
                [P(R(name, F_BOLD, 14, INK), align=PP_ALIGN.CENTER)])
        textbox(s, x, y + 1.06, tw, 0.32,
                [P(R(cap, F_MED, 10.5, GRAY), align=PP_ALIGN.CENTER)])
    # arrow down
    connector(s, SW / 2, 4.1, SW / 2, 4.72, color=GRAY2, weight=2.0, dash="dash")
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(SW/2 - 0.13),
                             Inches(4.66), Inches(0.26), Inches(0.2))
    tri.rotation = 180; tri.fill.solid(); tri.fill.fore_color.rgb = hx(GRAY2)
    tri.line.fill.background(); tri.shadow.inherit = False
    # result bar
    rb = card(s, MX, 5.0, SW - 2 * MX, 1.35, fill=INK, line=None, radius=0.2,
              shadow=True)
    textbox(s, MX + 0.5, 5.18, 7.6, 1.0,
            [P([R("Когнитивная перегрузка", F_XBOLD, 24, PAPER)], after=2),
             P(R("слишком много решений и источников информации за один вечер",
                 F_MED, 13.5, "C9C7E6"))],
            anchor=MSO_ANCHOR.MIDDLE)
    chip(s, SW - MX - 3.3, 5.46, 2.9, 0.62, "5+ сервисов · 0 структуры",
         INDIGO, PAPER, size=13)
    footer(s, 3)


# ============================================================== SLIDE 4 ======
def s04(prs):
    s = new(prs); base(s)
    header(s, "ИССЛЕДОВАНИЯ ЦА", "Что мы узнали из интервью", 4, title_size=32)
    textbox(s, MX, 1.92, 11.6, 0.4,
            [P(R("Глубинные интервью с ЦА — разные города, возраст и форматы. "
                 "Боли повторялись у всех.", F_MED, 14, GRAY))])
    cards = [
        (PINK,   "Не знают, чем заняться",
         "Хотят попробовать новое, но устают от поиска и откладывают."),
        (INDIGO, "Выбор съедает вечер",
         "30–90 минут на 4+ сервиса — и снова привычное место."),
        (GREEN,  "Нужен бюджет заранее",
         "Хотят понимать стоимость вечера ещё до того, как собрались."),
    ]
    gap = 0.3; cw = (SW - 2 * MX - 2 * gap) / 3; y = 2.5; h = 2.7
    for i, (col, title, body) in enumerate(cards):
        x = MX + i * (cw + gap)
        card(s, x, y, cw, h, radius=0.18, shadow=True)
        icon_tile(s, x + 0.34, y + 0.34, 0.6, col, kind="dot")
        textbox(s, x + 0.34, y + 1.12, cw - 0.6, 0.66,
                [P(R(title, F_XBOLD, 17, INK), ls=1.04)])
        textbox(s, x + 0.34, y + 1.84, cw - 0.62, 0.78,
                [P(R(body, F_MED, 13, GRAY), ls=1.14)])
    # highlight bar
    hb = card(s, MX, 5.45, SW - 2 * MX, 1.2, fill=LAV, line=None, radius=0.2,
              shadow=False)
    oval(s, MX + 0.42, 5.83, 0.42, 0.42, INDIGO)
    textbox(s, MX + 1.1, 5.55, SW - 2 * MX - 1.4, 1.0,
            [P([R("Главный вывод:  ", F_XBOLD, 16, INDIGO),
                R("понятного бюджета до встречи нет нигде в структурированном "
                  "виде — это и есть незакрытая боль.", F_MED, 15, INK)],
               ls=1.12)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 4)


# ============================================================== SLIDE 5 ======
def s05(prs):
    s = new(prs); base(s)
    header(s, "ЦЕЛЕВАЯ АУДИТОРИЯ", "Для кого эта проблема", 5, title_size=32)
    # left persona card
    px, pw, py, ph = MX, 5.25, 2.2, 4.0
    card(s, px, py, pw, ph, radius=0.2, shadow=True)
    oval(s, px + 0.45, py + 0.5, 1.1, 1.1, LAV2)
    oval(s, px + 0.62, py + 0.62, 0.76, 0.76, INDIGO)
    textbox(s, px + 0.45, py + 0.55, 1.1, 1.0,
            [P(R("18\n30", F_XBOLD, 17, PAPER), align=PP_ALIGN.CENTER, ls=0.95)],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, px + 1.78, py + 0.52, pw - 2.0, 0.5,
            [P(R("Молодые горожане", F_XBOLD, 20, INK))])
    textbox(s, px + 1.78, py + 0.98, pw - 2.0, 0.5,
            [P(R("18–30 лет · города от 300 000", F_MED, 13.5, INDIGO))])
    hline(s, px + 0.45, py + 1.95, pw - 0.9, color=LINE)
    textbox(s, px + 0.45, py + 2.16, pw - 0.9, 1.6,
            [P(R("Регулярно проводят досуг вне дома, ограничены в бюджете и "
                 "принимают решение сами или небольшой компанией.",
                 F_MED, 14, GRAY), ls=1.2)])
    # right: SAM stat + needs
    rx = px + pw + 0.5; rw = SW - MX - rx
    st = card(s, rx, py, rw, 1.5, fill=INK, line=None, radius=0.2, shadow=True)
    textbox(s, rx + 0.5, py + 0.26, rw - 1.0, 0.4,
            [P(R("РАЗМЕР АУДИТОРИИ (SAM)", F_BOLD, 11.5, INDIGO_L, sp=2.4))])
    textbox(s, rx + 0.46, py + 0.55, rw - 1.0, 0.9,
            [P([R("≈ 8 млн", F_BLACK, 40, PAPER),
                R("   человек\n   в 26 городах", F_MED, 14, "C9C7E6")], ls=1.05)],
            anchor=MSO_ANCHOR.MIDDLE)
    needs = [
        ("Быстро выбрать вариант", "без часа поиска по сервисам"),
        ("Понимать бюджет заранее", "сколько выйдет на человека"),
        ("Не терять вечер", "решение — за минуты, а не за час"),
    ]
    ny = py + 1.75; nh = 0.68; ng = 0.13
    for i, (t, sub) in enumerate(needs):
        y = ny + i * (nh + ng)
        c = card(s, rx, y, rw, nh, radius=0.14, shadow=False, line=LINE)
        icon_tile(s, rx + 0.22, y + 0.13, 0.42, GREEN, kind="check")
        textbox(s, rx + 0.82, y + 0.05, rw - 1.0, nh,
                [P([R(t + "  ", F_BOLD, 14.5, INK),
                    R("— " + sub, F_MED, 12.5, GRAY)])],
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 5)


# ============================================================== SLIDE 6 ======
def s06(prs):
    s = new(prs); base(s, INK); blobs_dark(s)
    textbox(s, MX, 0.95, 8.0, 0.4,
            [P(R("А ТЕПЕРЬ — ТА ЖЕ ПЯТНИЦА, ТА ЖЕ КОМПАНИЯ", F_BOLD, 13,
                 INDIGO_L, sp=2.6))])
    textbox(s, MX - 0.02, 1.42, 8.2, 1.7,
            [P([R("А если бы у компании\nбыл ", F_XBOLD, 38, PAPER),
                R("SpenDay", F_XBOLD, 38, INDIGO_L)], ls=1.06)])
    steps = [
        ("1", "Открываем ленту сценариев", "вместо четырёх разных сервисов"),
        ("2", "Ставим фильтры", "«с друзьями» · город · бюджет · формат"),
        ("3", "Видим готовые варианты", "с бюджетом на человека сразу"),
        ("4", "Один понятный путь", "вместо хаоса из карт, соцсетей и чатов"),
    ]
    y0 = 3.35; rh = 0.84
    for i, (n, t, sub) in enumerate(steps):
        y = y0 + i * rh
        oval(s, MX, y, 0.5, 0.5, INDIGO)
        textbox(s, MX, y, 0.5, 0.5,
                [P(R(n, F_XBOLD, 18, PAPER), align=PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, MX + 0.72, y - 0.04, 6.2, 0.4,
                [P(R(t, F_BOLD, 16.5, PAPER))])
        textbox(s, MX + 0.72, y + 0.3, 6.4, 0.4,
                [P(R(sub, F_MED, 12.5, "B9B7DE"))])
    glow(s, 10.7, 4.0, 4.5, 6.0, GLOW_D)
    phone(s, "feed", height=6.55, cx=10.75, top=1.05)
    footer(s, 6, on_dark=True)


# --------- product slide scaffold: phone(s) + side panel ---------------------
def attr_row(s, x, y, w, color, title, body):
    oval(s, x, y + 0.06, 0.28, 0.28, color)
    textbox(s, x + 0.48, y - 0.04, w - 0.5, 0.4,
            [P([R(title, F_BOLD, 15, INK)])])
    textbox(s, x + 0.48, y + 0.3, w - 0.5, 0.4,
            [P(R(body, F_MED, 12.5, GRAY))])


# ============================================================== SLIDE 7 ======
def s07(prs):
    s = new(prs); base(s)
    header(s, "ПРОДУКТ · UX", "Выбор готового сценария", 7, title_size=32)
    glow(s, 3.2, 4.5, 4.0, 5.6, GLOW)
    phone(s, "feed2", height=5.0, cx=3.2, top=1.95)
    rx, rw = 6.0, SW - MX - 6.0
    textbox(s, rx, 2.05, rw, 0.5,
            [P(R("Карточка показывает всё, что нужно для решения — сразу:",
                 F_BOLD, 16, INK), ls=1.15)])
    rows = [
        (PINK,   "Название и ситуация", "свидание, с друзьями, прогулка…"),
        (INDIGO, "Бюджет на человека", "понятная сумма ещё до клика"),
        (INK,    "Длительность и места", "сколько времени и где"),
        (GREEN,  "Формат и теги", "дома, в городе, на природе"),
    ]
    y = 2.95
    for col, t, b in rows:
        attr_row(s, rx, y, rw, col, t, b)
        y += 0.82
    hb = card(s, rx, y + 0.05, rw, 0.85, fill=LAV, line=None, radius=0.16)
    textbox(s, rx + 0.32, y + 0.05, rw - 0.6, 0.85,
            [P([R("Не список мест — ", F_BOLD, 14.5, INDIGO),
                R("готовый вариант вечера.", F_MED, 14.5, INK)])],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 7)


# ============================================================== SLIDE 8 ======
def s08(prs):
    s = new(prs); base(s)
    header(s, "ПРОДУКТ · UX", "Просмотр сценария", 8, title_size=32)
    glow(s, 3.2, 4.5, 4.0, 5.6, GLOW)
    phone(s, "scenario", height=5.05, cx=3.2, top=1.9)
    rx, rw = 6.0, SW - MX - 6.0
    textbox(s, rx, 2.0, rw, 0.5,
            [P(R("Внутри сценария — весь вечер по шагам и стоимости:",
                 F_BOLD, 16, INK), ls=1.15)])
    steps = [
        ("1", "Прогулка по набережной", "бесплатно", GREEN),
        ("2", "Ужин в ресторане", "от 2 500 ₽", INK),
        ("3", "Коктейли в баре", "от 1 700 ₽", INK),
    ]
    y = 2.78; sh = 0.74
    for n, t, price, pc in steps:
        c = card(s, rx, y, rw, sh - 0.1, radius=0.14, shadow=False, line=LINE)
        oval(s, rx + 0.2, y + 0.13, 0.38, 0.38, INDIGO)
        textbox(s, rx + 0.2, y + 0.13, 0.38, 0.38,
                [P(R(n, F_BOLD, 13, PAPER), align=PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, rx + 0.74, y + 0.02, rw - 2.1, sh - 0.1,
                [P(R(t, F_BOLD, 14, INK))], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, rx + rw - 1.5, y + 0.02, 1.35, sh - 0.1,
                [P(R(price, F_BOLD, 13, pc), align=PP_ALIGN.RIGHT)],
                anchor=MSO_ANCHOR.MIDDLE)
        y += sh
    y += 0.08
    chip(s, rx, y, 2.55, 0.66, "~ 2 100 ₽ на человека", INDIGO, PAPER, size=13.5)
    rect(s, rx + 2.7, y, 2.45, 0.66, fill=LAV, radius=0.33)
    star(s, rx + 2.96, y + 0.2, 0.27, AMBER)
    textbox(s, rx + 3.33, y, 1.75, 0.66,
            [P(R("4,8 · 124 отзыва", F_BOLD, 13.5, INK))],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, rx, y + 0.92, rw, 0.5,
            [P(R("Пользователь сразу понимает, что будет и сколько это стоит.",
                 F_MED, 13.5, GRAY), ls=1.15)])
    footer(s, 8)


# ============================================================== SLIDE 9 ======
def s09(prs):
    s = new(prs); base(s)
    header(s, "ПРОДУКТ · UX", "Согласование с компанией", 9, title_size=32)
    glow(s, 4.0, 4.6, 6.6, 5.6, GLOW)
    phone(s, "vote", height=4.85, cx=2.85, top=2.1)
    phone(s, "vote2", height=4.85, cx=5.5, top=2.1)
    # arrow between phones
    rx, rw = 7.7, SW - MX - 7.7
    textbox(s, rx, 2.05, rw, 0.5,
            [P(R("Скинул ссылку — компания выбрала за 5 минут:",
                 F_BOLD, 16.5, INK), ls=1.15)])
    pts = [
        ("Поделиться сценариями в общий чат", INDIGO),
        ("Друзья голосуют без регистрации и скачивания", GREEN),
        ("Результат виден всем в реальном времени", PINK),
    ]
    y = 3.0
    for t, col in pts:
        icon_tile(s, rx, y, 0.46, col, kind="check")
        textbox(s, rx + 0.66, y - 0.04, rw - 0.7, 0.6,
                [P(R(t, F_MED, 14.5, INK), ls=1.1)], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.86
    hb = card(s, rx, y + 0.06, rw, 0.95, fill=INK, line=None, radius=0.18,
              shadow=True)
    textbox(s, rx + 0.34, y + 0.06, rw - 0.6, 0.95,
            [P([R("Канал роста: ", F_BOLD, 14, INDIGO_L),
                R("каждое голосование — это 3–7 человек, которые видят SpenDay.",
                  F_MED, 13.5, PAPER)], ls=1.15)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 9)


# ============================================================== SLIDE 10 =====
def s10(prs):
    s = new(prs); base(s)
    header(s, "ПРОДУКТ · UX", "После встречи: расходы и сохранение", 10,
           title_size=30)
    glow(s, 4.0, 4.6, 6.6, 5.6, GLOW)
    phone(s, "calc", height=4.8, cx=2.85, top=2.05)
    phone(s, "profile", height=4.8, cx=5.5, top=2.05)
    rx, rw = 7.7, SW - MX - 7.7
    textbox(s, rx, 2.0, rw, 0.5,
            [P(R("SpenDay закрывает весь путь вечера:", F_BOLD, 16.5, INK))])
    pts = [
        ("Калькулятор расходов", "кто кому сколько должен — без регистрации"),
        ("Сохранение сценария", "удачный вечер — в один тап"),
        ("Личный календарь", "история вечеров и план на будущее"),
    ]
    y = 2.7
    for t, b in pts:
        icon_tile(s, rx, y, 0.46, INDIGO, kind="check")
        textbox(s, rx + 0.66, y - 0.06, rw - 0.7, 0.66,
                [P([R(t + "  ", F_BOLD, 14.5, INK)], after=1),
                 P(R(b, F_MED, 12, GRAY))])
        y += 0.84
    # loop strip
    loop = ["Выбрать", "Согласовать", "Сходить", "Разделить", "Сохранить"]
    ly = y + 0.12; lx = rx
    cw = (rw - 0.0) / len(loop)
    strip = card(s, rx, ly, rw, 0.62, fill=LAV, line=None, radius=0.16)
    for i, w in enumerate(loop):
        textbox(s, rx + i * cw, ly, cw, 0.62,
                [P(R(w, F_BOLD, 11.5, INDIGO if i == 0 or i == len(loop)-1 else INK),
                   align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        if i < len(loop) - 1:
            textbox(s, rx + (i + 1) * cw - 0.12, ly, 0.24, 0.62,
                    [P(R("›", F_BOLD, 14, GRAY2), align=PP_ALIGN.CENTER)],
                    anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 10)


# ============================================================== SLIDE 11 =====
def s11(prs):
    s = new(prs); base(s)
    header(s, "РЫНОК И КОНКУРЕНТЫ", "Какие альтернативы есть сейчас", 11,
           title_size=32)
    textbox(s, MX, 1.92, 11.6, 0.4,
            [P(R("Каждый сервис закрывает свой кусок пути — собрать решение "
                 "всё равно приходится самому.", F_MED, 14, GRAY))])
    cards = [
        ("Соцсети", "TikTok, Instagram, Pinterest", "Дают идеи и вдохновение",
         "Нет плана, цен и согласования", PINK),
        ("Карты", "2ГИС, Яндекс.Карты", "Находят место и маршрут",
         "Нужно уже знать, что ищешь", INDIGO),
        ("Афиши", "Афиша, Tripadvisor", "События и билеты",
         "Дорогой и редкий сегмент", AMBER),
        ("Деление расходов", "Splitwise, Tricount", "Делят счёт после траты",
         "Не отвечают «куда идём»", GREEN),
    ]
    gap = 0.28; cw = (SW - 2 * MX - 3 * gap) / 4; y = 2.6; h = 3.7
    for i, (name, brands, good, gap_t, col) in enumerate(cards):
        x = MX + i * (cw + gap)
        card(s, x, y, cw, h, radius=0.18, shadow=True)
        rect(s, x, y, cw, 0.14, fill=col, radius=0.07)
        textbox(s, x + 0.3, y + 0.4, cw - 0.55, 0.5,
                [P(R(name, F_XBOLD, 16.5, INK))])
        textbox(s, x + 0.3, y + 0.86, cw - 0.55, 0.5,
                [P(R(brands, F_MED, 11, GRAY2))])
        hline(s, x + 0.3, y + 1.42, cw - 0.6, color=LINE)
        textbox(s, x + 0.3, y + 1.62, cw - 0.55, 0.3,
                [P(R("ЧТО ДАЮТ", F_BOLD, 9.5, GREEN_D, sp=1.6))])
        textbox(s, x + 0.3, y + 1.9, cw - 0.55, 0.8,
                [P(R(good, F_MED, 13, INK), ls=1.1)])
        textbox(s, x + 0.3, y + 2.72, cw - 0.55, 0.3,
                [P(R("ЧЕГО НЕ ХВАТАЕТ", F_BOLD, 9.5, PINK, sp=1.4))])
        textbox(s, x + 0.3, y + 3.0, cw - 0.55, 0.8,
                [P(R(gap_t, F_MED, 13, INK), ls=1.1)])
    footer(s, 11)


# ============================================================== SLIDE 12 =====
def s12(prs):
    s = new(prs); base(s)
    header(s, "РЫНОК И КОНКУРЕНТЫ",
           "Почему конкуренты не закрывают путь целиком", 12, title_size=28)
    cols = ["SpenDay", "Соцсети", "Splitwise", "Карты", "Афиша"]
    rows = [
        ("Готовый сценарий под ситуацию", [1, 0, 0, 0, 0]),
        ("Бюджет известен заранее",       [1, 0, 0, 0, 1]),
        ("Групповое согласование",        [1, 0, 1, 0, 0]),
        ("Базовые функции без регистрации",[1, 1, 0, 1, 0]),
        ("Повседневный досуг",            [1, 1, 0, 1, 0]),
    ]
    tx, ty = MX, 2.15
    label_w = 4.3
    col_w = (SW - 2 * MX - label_w) / len(cols)
    rh = 0.62
    # header row
    for j, c in enumerate(cols):
        cx = tx + label_w + j * col_w
        hl = (j == 0)
        if hl:
            rect(s, cx + 0.06, ty - 0.05, col_w - 0.12, rh + 0.04, fill=INDIGO,
                 radius=0.12)
        textbox(s, cx, ty, col_w, rh,
                [P(R(c, F_XBOLD if hl else F_BOLD, 13 if hl else 12.5,
                     PAPER if hl else INK), align=PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
    # body rows
    for i, (label, vals) in enumerate(rows):
        y = ty + rh + i * rh
        if i % 2 == 0:
            rect(s, tx, y, label_w, rh, fill=BG)
        # SpenDay column band
        rect(s, tx + label_w + 0.06, y, col_w - 0.12, rh,
             fill="EEEDFB")
        textbox(s, tx + 0.18, y, label_w - 0.3, rh,
                [P(R(label, F_MED, 13, INK))], anchor=MSO_ANCHOR.MIDDLE)
        for j, v in enumerate(vals):
            cx = tx + label_w + j * col_w
            if v:
                col = INDIGO if j == 0 else GREEN
                oval(s, cx + col_w/2 - 0.16, y + rh/2 - 0.16, 0.32, 0.32, col)
                icon_tile(s, cx + col_w/2 - 0.16, y + rh/2 - 0.16, 0.32, col,
                          kind="check")
            else:
                textbox(s, cx, y, col_w, rh,
                        [P(R("—", F_BOLD, 15, "C5C5D6"), align=PP_ALIGN.CENTER)],
                        anchor=MSO_ANCHOR.MIDDLE)
    # journey strip
    jy = ty + rh * (len(rows) + 1) + 0.18
    journey = ["Ситуация", "Сценарий", "Бюджет", "Согласование", "Расходы"]
    jb = card(s, MX, jy, SW - 2*MX, 0.95, fill=INK, line=None, radius=0.18,
              shadow=True)
    textbox(s, MX + 0.34, jy + 0.12, 3.0, 0.7,
            [P(R("Один путь\nв SpenDay:", F_BOLD, 13, INDIGO_L), ls=1.05)],
            anchor=MSO_ANCHOR.MIDDLE)
    jx = MX + 3.0; jw = (SW - 2*MX) - 3.3
    step_w = jw / len(journey)
    for i, st in enumerate(journey):
        textbox(s, jx + i * step_w, jy, step_w, 0.95,
                [P(R(st, F_BOLD, 13, PAPER), align=PP_ALIGN.CENTER)],
                anchor=MSO_ANCHOR.MIDDLE)
        if i < len(journey) - 1:
            textbox(s, jx + (i+1)*step_w - 0.14, jy, 0.28, 0.95,
                    [P(R("→", F_BOLD, 15, INDIGO_L), align=PP_ALIGN.CENTER)],
                    anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 12)


# ============================================================== SLIDE 13 =====
def s13(prs):
    s = new(prs); base(s)
    header(s, "УНИКАЛЬНОЕ ПРЕДЛОЖЕНИЕ", "УТП SpenDay", 13, title_size=32)
    # statement panel
    sp = card(s, MX, 2.1, SW - 2*MX, 1.85, fill=INK, line=None, radius=0.22,
              shadow=True)
    textbox(s, MX + 0.55, 2.32, SW - 2*MX - 1.1, 0.4,
            [P(R("ГЛАВНАЯ МЫСЛЬ", F_BOLD, 11.5, INDIGO_L, sp=2.8))])
    textbox(s, MX + 0.5, 2.66, SW - 2*MX - 1.0, 1.2,
            [P([R("SpenDay продаёт не места, а ", F_XBOLD, 23, PAPER),
                R("быстрое решение", F_XBOLD, 23, INDIGO_L),
                R(": готовый сценарий под ситуацию с понятным бюджетом — "
                  "выбрать, согласовать и разделить расходы в одном приложении.",
                  F_MED, 19, "D8D6F0")], ls=1.16)])
    pillars = [
        ("Сценарий под ситуацию", "готовый план вечера, а не список мест", PINK),
        ("Понятный бюджет заранее", "стоимость на человека до клика", INDIGO),
        ("Согласование без регистрации", "голосование по ссылке в чате", GREEN),
        ("Калькулятор расходов", "кто кому сколько — за секунды", AMBER),
    ]
    gap = 0.28; cw = (SW - 2*MX - 3*gap)/4; y = 4.35; h = 2.25
    for i, (t, b, col) in enumerate(pillars):
        x = MX + i*(cw+gap)
        card(s, x, y, cw, h, radius=0.18, shadow=True)
        icon_tile(s, x + 0.3, y + 0.32, 0.56, col, kind="dot")
        textbox(s, x + 0.3, y + 1.04, cw - 0.55, 0.6,
                [P(R(t, F_XBOLD, 15, INK), ls=1.05)])
        textbox(s, x + 0.3, y + 1.6, cw - 0.55, 0.6,
                [P(R(b, F_MED, 12.5, GRAY), ls=1.12)])
    footer(s, 13)


# ============================================================== SLIDE 14 =====
def s14(prs):
    s = new(prs); base(s)
    header(s, "МОНЕТИЗАЦИЯ", "Как SpenDay зарабатывает", 14, title_size=32)
    textbox(s, MX, 1.92, 11.6, 0.4,
            [P([R("Базовый продукт бесплатный", F_BOLD, 14, INDIGO),
                R(" — это растит аудиторию и не мешает UX. Доход подключается "
                  "поэтапно.", F_MED, 14, GRAY)])])
    stages = [
        ("Старт", "0–10 000 MAU",
         ["Реклама РСЯ (eCPM 50–110 ₽)", "Нативные партнёрства 5–20k ₽/мес"],
         "ARPU ~2–3 ₽", INDIGO),
        ("Рост", "10 000–100 000 MAU",
         ["+ Платный UGC: комиссия с авторов", "Подписка на авторов сценариев"],
         "ARPU ~5 ₽", PINK),
        ("Масштаб", "100 000+ MAU",
         ["+ Агрегация и комиссия с покупок", "Bundle: билеты, брони, доставка"],
         "ARPU ~8 ₽", GREEN),
    ]
    gap = 0.3; cw = (SW - 2*MX - 2*gap)/3; y = 2.6; h = 2.95
    for i, (name, mau, items, arpu, col) in enumerate(stages):
        x = MX + i*(cw+gap)
        card(s, x, y, cw, h, radius=0.18, shadow=True)
        chip(s, x + 0.3, y + 0.3, 1.55, 0.46, name, col, PAPER, size=12.5)
        textbox(s, x + 0.32, y + 0.86, cw - 0.6, 0.32,
                [P(R(mau, F_BOLD, 11.5, GRAY))])
        yy = y + 1.34
        for it in items:
            oval(s, x + 0.32, yy + 0.07, 0.12, 0.12, col)
            textbox(s, x + 0.56, yy - 0.04, cw - 0.85, 0.5,
                    [P(R(it, F_MED, 12.5, INK), ls=1.05)])
            yy += 0.54
        hline(s, x + 0.3, y + h - 0.6, cw - 0.6, color=LINE)
        textbox(s, x + 0.3, y + h - 0.52, cw - 0.6, 0.42,
                [P(R(arpu, F_XBOLD, 16, col))], anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            textbox(s, x + cw + 0.02, y + 1.25, 0.26, 0.5,
                    [P(R("›", F_BLACK, 20, GRAY2), align=PP_ALIGN.CENTER)])
    # break-even strip
    by = y + h + 0.2
    bb = card(s, MX, by, SW - 2*MX, 0.92, fill=INK, line=None, radius=0.18,
              shadow=True)
    textbox(s, MX + 0.4, by, 6.5, 0.92,
            [P([R("Точка безубыточности:  ", F_MED, 15, "C9C7E6"),
                R("10 000 MAU", F_XBOLD, 18, PAPER)])],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MX + 6.6, by, SW - 2*MX - 7.0, 0.92,
            [P(R("20 000 ₽ серверы/мес  ÷  2 ₽ ARPU.  Дальше — выручка кратно "
                 "обгоняет расходы.", F_MED, 12.5, "B9B7DE"), ls=1.1)],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 14)


# ============================================================== SLIDE 15 =====
def s15(prs):
    s = new(prs); base(s)
    header(s, "MVP", "Что входит в MVP", 15, title_size=32)
    textbox(s, MX, 1.92, 11.6, 0.4,
            [P(R("Минимум, чтобы проверить главную гипотезу — и собрать его "
                 "можно быстро.", F_MED, 14, GRAY))])
    feats = [
        ("Лента сценариев", "готовые варианты вечера"),
        ("Фильтры", "ситуация · бюджет · формат"),
        ("Карточка сценария", "шаги, бюджет, длительность"),
        ("Шаринг ссылкой", "отправить компании в чат"),
        ("Голосование", "без регистрации и скачивания"),
        ("Калькулятор расходов", "разделить счёт за секунды"),
    ]
    gap = 0.3; cw = (SW - 2*MX - 2*gap)/3; rh = 1.4; y0 = 2.55
    for i, (t, b) in enumerate(feats):
        r, cidx = divmod(i, 3)
        x = MX + cidx*(cw+gap); y = y0 + r*(rh+0.25)
        card(s, x, y, cw, rh, radius=0.16, shadow=True)
        icon_tile(s, x + 0.3, y + 0.3, 0.5, GREEN, kind="check")
        textbox(s, x + 0.98, y + 0.26, cw - 1.1, 0.45,
                [P(R(t, F_XBOLD, 15.5, INK))])
        textbox(s, x + 0.98, y + 0.72, cw - 1.1, 0.5,
                [P(R(b, F_MED, 12.5, GRAY), ls=1.05)])
    by = y0 + 2*(rh) + 0.25 + 0.1
    bb = card(s, MX, by, SW - 2*MX, 0.92, fill=LAV, line=None, radius=0.18)
    textbox(s, MX + 0.4, by, SW - 2*MX - 0.8, 0.92,
            [P([R("Цель MVP:  ", F_XBOLD, 15.5, INDIGO),
                R("проверить, что люди выбирают досуг быстрее и возвращаются за "
                  "новыми сценариями.", F_MED, 15, INK)], ls=1.12)],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 15)


# ============================================================== SLIDE 16 =====
def s16(prs):
    s = new(prs); base(s)
    header(s, "МЕТРИКИ", "Как измеряем успех", 16, title_size=32)
    # North star
    ns = card(s, MX, 2.1, SW - 2*MX, 1.25, fill=INK, line=None, radius=0.2,
              shadow=True)
    textbox(s, MX + 0.5, 2.28, 3.4, 0.4,
            [P(R("NORTH STAR METRIC", F_BOLD, 11.5, INDIGO_L, sp=2.6))])
    textbox(s, MX + 0.46, 2.58, 8.5, 0.7,
            [P([R("MAU", F_BLACK, 26, PAPER),
                R("  +  ", F_MED, 22, INDIGO_L),
                R("Monthly Retention", F_BLACK, 26, PAPER)])],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, SW - MX - 4.3, 2.5, 4.0, 0.85,
            [P(R("Масштаб аудитории + качество роста — самая честная пара "
                 "для продукта планирования досуга.", F_MED, 12, "C9C7E6"),
               ls=1.15)], anchor=MSO_ANCHOR.MIDDLE)
    metrics = [
        ("Time to Save", "до 3 мин", "vs 30–90 мин сейчас", INDIGO),
        ("Monthly Retention", "25–35%", "ниже 15% — сигнал", PINK),
        ("DAU / MAU", "12–18%", "норма lifestyle", GREEN),
        ("Сохранений / польз.", "5–15", "ценность контента", AMBER),
        ("Участники голосования", "4+", "виральность", INDIGO),
        ("CAC", "≤ 100 ₽", "LTV/CAC > 2x", PINK),
    ]
    gap = 0.28; cw = (SW - 2*MX - 2*gap)/3; rh = 1.32; y0 = 3.7
    for i, (name, val, sub, col) in enumerate(metrics):
        r, cidx = divmod(i, 3)
        x = MX + cidx*(cw+gap); y = y0 + r*(rh+0.22)
        card(s, x, y, cw, rh, radius=0.16, shadow=True)
        rect(s, x, y + 0.22, 0.1, rh - 0.44, fill=col, radius=0.05)
        textbox(s, x + 0.34, y + 0.22, cw - 0.5, 0.35,
                [P(R(name, F_BOLD, 12.5, GRAY))])
        textbox(s, x + 0.32, y + 0.5, cw - 0.5, 0.55,
                [P(R(val, F_BLACK, 25, INK))])
        textbox(s, x + 0.34, y + 1.02, cw - 0.5, 0.3,
                [P(R(sub, F_MED, 11, col))])
    footer(s, 16)


# ============================================================== SLIDE 17 =====
def s17(prs):
    s = new(prs); base(s, INK); blobs_dark(s)
    wordmark(s, MX, 0.8, size=0.4, on_dark=True)
    textbox(s, MX - 0.02, 1.75, 8.6, 1.5,
            [P([R("Почему в SpenDay\nможно ", F_XBOLD, 36, PAPER),
                R("инвестировать", F_XBOLD, 36, INDIGO_L)], ls=1.05)])
    reasons = [
        ("Боль массовая и ежедневная", "SAM ≈ 8 млн человек, 18–30 лет в городах"),
        ("Запуск дёшев", "MVP на простом стеке, безубыточность с 10 000 MAU"),
        ("Рост встроен в продукт", "голосование без регистрации = viral-петля"),
        ("Понятная экономика", "ARPU растёт 2 → 8 ₽, модель прибыльна"),
    ]
    y = 3.45
    for t, b in reasons:
        icon_tile(s, MX, y, 0.5, INDIGO, kind="check")
        textbox(s, MX + 0.72, y - 0.06, 7.4, 0.66,
                [P([R(t + "  ", F_BOLD, 16.5, PAPER)], after=1),
                 P(R(b, F_MED, 12.5, "B9B7DE"))])
        y += 0.86
    # CTA panel right
    cx, cw = 8.9, SW - MX - 8.9
    cta = card(s, cx, 3.45, cw, 3.45, fill=INK2, line=INK3, lw=1.2, radius=0.22,
               shadow=False)
    textbox(s, cx + 0.4, 3.75, cw - 0.8, 0.4,
            [P(R("ИНВЕСТИЦИИ ИДУТ В УСКОРЕНИЕ", F_BOLD, 11, INDIGO_L, sp=2.0))])
    for i, t in enumerate(["Контент и сценарии", "Микро-креаторы",
                           "Продвижение", "Партнёрства"]):
        yy = 4.18 + i*0.5
        oval(s, cx + 0.42, yy + 0.07, 0.14, 0.14, INDIGO_L)
        textbox(s, cx + 0.72, yy - 0.04, cw - 1.0, 0.4,
                [P(R(t, F_MED, 14, PAPER))])
    hline(s, cx + 0.4, 6.3, cw - 0.8, color=INK3)
    textbox(s, cx + 0.4, 6.42, cw - 0.8, 0.45,
            [P(R("Выбрать вечер за 10 минут — а не за вечер.",
                 F_BOLD, 13.5, INDIGO_L), ls=1.1)])
    footer(s, 17, on_dark=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
               s11, s12, s13, s14, s15, s16, s17]:
        fn(prs)
    prs.save(OUT)
    print("saved", OUT, "·", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    build()
