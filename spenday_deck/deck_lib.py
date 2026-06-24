"""Reusable design-system helpers for the SpenDay pitch deck (python-pptx)."""
import os
from lxml import etree
from PIL import Image as PImage
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ASSETS = os.path.join(os.path.dirname(__file__), "assets")

# ---------------------------------------------------------------- palette ----
INK      = "1E1B4A"   # near-black navy (headlines, dark backgrounds)
INK2     = "2A2660"   # raised panel on navy
INK3     = "3A357A"   # hairline on navy
INDIGO   = "5046E5"   # brand primary
INDIGO_D = "3F37C9"
INDIGO_L = "8B83FF"   # light indigo for use on dark backgrounds
GREEN    = "14B886"   # money / free / positive
GREEN_D  = "0E9A6F"
PINK     = "DB4E8F"   # situations / energy
AMBER    = "F5A524"   # rating stars
PAPER    = "FFFFFF"
BG       = "F6F6FC"   # light page tint
LAV      = "EFEEF9"   # light lavender panel
LAV2     = "E7E6F6"   # device tray / soft fill
GLOW     = "ECEBFB"   # phone glow on light
GLOW_D   = "2A2766"   # decorative blob on navy
GRAY     = "6A6A86"   # muted body text
GRAY2    = "9A9AB2"   # lighter caption
LINE     = "E6E6F1"   # hairline on light

# ------------------------------------------------------------------- fonts ----
F_BLACK = "Montserrat Black"
F_XBOLD = "Montserrat ExtraBold"
F_BOLD  = "Montserrat SemiBold"
F_MED   = "Montserrat Medium"
F_REG   = "Montserrat"
F_LIGHT = "Montserrat Light"

EMU_IN = 914400
PHONE_AR = {}  # cache of png aspect ratios


def hx(c):
    return RGBColor.from_string(c)


def _in(v):
    return Inches(v) if not isinstance(v, Emu) else v


# ----------------------------------------------------------- text plumbing ----
def R(t, f=F_REG, s=14, c=INK, sp=None, italic=False):
    """run-spec"""
    return dict(t=t, f=f, s=s, c=c, sp=sp, italic=italic)


def P(runs, align=PP_ALIGN.LEFT, before=0, after=0, ls=None):
    """paragraph-spec"""
    if isinstance(runs, dict):
        runs = [runs]
    return dict(runs=runs, align=align, before=before, after=after, ls=ls)


def _spc(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def _fill_tf(tf, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        if pa.get("before"):
            p.space_before = Pt(pa["before"])
        p.space_after = Pt(pa.get("after", 0) or 0)
        if pa.get("ls"):
            p.line_spacing = pa["ls"]
        for rs in pa["runs"]:
            r = p.add_run()
            r.text = rs["t"]
            r.font.name = rs["f"]
            r.font.size = Pt(rs["s"])
            r.font.color.rgb = hx(rs["c"])
            r.font.italic = rs.get("italic", False)
            if rs.get("sp") is not None:
                _spc(r, rs["sp"])


def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    _fill_tf(tb.text_frame, paras, anchor, wrap)
    return tb


def set_text(shape, paras, anchor=MSO_ANCHOR.MIDDLE, wrap=True, pad=0.0):
    tf = shape.text_frame
    if pad:
        tf.margin_left = tf.margin_right = Inches(pad)
    _fill_tf(tf, paras, anchor, wrap)
    if pad:
        tf.margin_left = tf.margin_right = Inches(pad)
    return shape


# ----------------------------------------------------------------- shapes ----
def _round(shape, radius_in):
    try:
        frac = Inches(radius_in) / min(int(shape.width), int(shape.height))
        shape.adjustments[0] = max(0.0, min(0.5, frac))
    except Exception:
        pass


def _soft_shadow(shape, blur=0.10, dist=0.055, alpha=20, color=INK, direction=5400000):
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    eff = etree.SubElement(spPr, qn("a:effectLst"))
    sh = etree.SubElement(eff, qn("a:outerShdw"))
    sh.set("blurRad", str(int(Inches(blur))))
    sh.set("dist", str(int(Inches(dist))))
    sh.set("dir", str(direction))
    sh.set("rotWithShape", "0")
    clr = etree.SubElement(sh, qn("a:srgbClr"))
    clr.set("val", color)
    a = etree.SubElement(clr, qn("a:alpha"))
    a.set("val", str(int(alpha * 1000)))


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, radius=None,
         shadow=False, dash=None):
    typ = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(typ, _in(x), _in(y), _in(w), _in(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = hx(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = hx(line)
        s.line.width = Pt(lw)
    if dash:
        ln = s.line._get_or_add_ln()
        d = etree.SubElement(ln, qn("a:prstDash"))
        d.set("val", dash)
    if radius:
        _round(s, radius)
    if shadow:
        _soft_shadow(s)
    return s


def oval(slide, x, y, w, h, fill, line=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, _in(x), _in(y), _in(w), _in(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = hx(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = hx(line)
        s.line.width = Pt(lw)
    return s


def glow(slide, cx, cy, w, h, color=GLOW):
    return oval(slide, cx - w / 2, cy - h / 2, w, h, color)


def hline(slide, x, y, w, color=LINE, weight=1.2):
    ln = slide.shapes.add_connector(2, _in(x), _in(y), _in(x + w), _in(y))
    ln.line.color.rgb = hx(color)
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def connector(slide, x1, y1, x2, y2, color=INDIGO, weight=1.6, dash=None):
    ln = slide.shapes.add_connector(2, _in(x1), _in(y1), _in(x2), _in(y2))
    ln.line.color.rgb = hx(color)
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    if dash:
        el = ln.line._get_or_add_ln()
        d = etree.SubElement(el, qn("a:prstDash"))
        d.set("val", dash)
    return ln


def chip(slide, x, y, w, h, text, fill, txt, size=11, font=F_BOLD,
         line=None, lw=1.0, sp=None, radius=None):
    s = rect(slide, x, y, w, h, fill=fill, line=line, lw=lw,
             radius=radius if radius is not None else h / 2)
    set_text(s, [P(R(text, font, size, txt, sp=sp), align=PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    return s


def star(slide, x, y, size, fill):
    try:
        s = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, _in(x), _in(y),
                                   _in(size), _in(size))
    except Exception:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, _in(x), _in(y),
                                   _in(size), _in(size))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = hx(fill)
    s.line.fill.background()
    return s


# ----------------------------------------------------------------- images ----
def _ar(name):
    if name not in PHONE_AR:
        iw, ih = PImage.open(os.path.join(ASSETS, f"phone_{name}.png")).size
        PHONE_AR[name] = iw / ih
    return PHONE_AR[name]


def phone(slide, name, height, left=None, top=0.0, cx=None, rot=0):
    path = os.path.join(ASSETS, f"phone_{name}.png")
    ar = _ar(name)
    h_emu = Inches(height)
    w_emu = int(h_emu * ar)
    if cx is not None:
        left_emu = int(Inches(cx) - w_emu / 2)
    else:
        left_emu = _in(left)
    pic = slide.shapes.add_picture(path, left_emu, _in(top), width=w_emu, height=h_emu)
    if rot:
        pic.rotation = rot
    return pic


# the visible phone (inside the padded transparent canvas) occupies ~ this
# fraction of the png height/width — useful when callers want the screen box
PHONE_PAD = 70  # px, must match process_assets.PAD


# ---------------------------------------------------------- slide scaffold ----
SW, SH = 13.333, 7.5
MX = 0.85  # left/right margin


def base(slide, color=PAPER):
    rect(slide, -0.06, -0.06, SW + 0.12, SH + 0.12, fill=color)


def wordmark(slide, x, y, size=0.40, on_dark=False, gap=0.14):
    """SpenDay lockup: rounded indigo square + star + 'SpenDay'."""
    sq = rect(slide, x, y, size, size, fill=INDIGO, radius=size * 0.28)
    star(slide, x + size * 0.27, y + size * 0.28, size * 0.46, PAPER)
    spen_c = PAPER if on_dark else INK
    day_c = INDIGO_L if on_dark else INDIGO
    fs = size * 52  # pt roughly proportional
    textbox(slide, x + size + gap, y - 0.06, 3.4, size + 0.12,
            [P([R("Spen", F_XBOLD, fs, spen_c), R("Day", F_XBOLD, fs, day_c)])],
            anchor=MSO_ANCHOR.MIDDLE)


def eyebrow(slide, text, x=MX, y=0.62, color=INDIGO):
    textbox(slide, x, y, 9.0, 0.32,
            [P(R(text, F_BOLD, 12.5, color, sp=2.4))])


def headline(slide, text, x=MX, y=0.98, w=11.4, size=33, color=INK, parts=None):
    runs = parts if parts else [R(text, F_XBOLD, size, color)]
    textbox(slide, x, y, w, 1.3, [P(runs, ls=1.04)])


def footer(slide, idx, total=17, on_dark=False):
    c1 = "C9C7E6" if on_dark else GRAY2
    c2 = INDIGO_L if on_dark else INDIGO
    cw = PAPER if on_dark else INK
    textbox(slide, MX, SH - 0.5, 3.0, 0.3,
            [P([R("Spen", F_BOLD, 10.5, cw), R("Day", F_BOLD, 10.5, c2)])],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, SW - MX - 2.0, SH - 0.5, 2.0, 0.3,
            [P([R(f"{idx:02d}", F_BOLD, 10.5, c2),
                R(f"  /  {total}", F_MED, 10.5, c1)], align=PP_ALIGN.RIGHT)],
            anchor=MSO_ANCHOR.MIDDLE)


def header(slide, eyebrow_text, title, idx, parts=None, title_size=33,
           title_w=11.4, on_dark=False):
    eyebrow(slide, eyebrow_text, color=(INDIGO_L if on_dark else INDIGO))
    headline(slide, title, size=title_size, w=title_w,
             color=PAPER if on_dark else INK, parts=parts)
    footer(slide, idx, on_dark=on_dark)


# tidy icon: rounded square tile with a glyph-ish mark (we draw simple marks)
def icon_tile(slide, x, y, size, fill, glyph_color=PAPER, kind="dot"):
    t = rect(slide, x, y, size, size, fill=fill, radius=size * 0.28)
    cx, cy = x + size / 2, y + size / 2
    if kind == "check":
        connector(slide, cx - size*0.18, cy + size*0.02, cx - size*0.04, cy + size*0.16,
                  color=glyph_color, weight=2.6)
        connector(slide, cx - size*0.04, cy + size*0.16, cx + size*0.20, cy - size*0.16,
                  color=glyph_color, weight=2.6)
    elif kind == "dot":
        oval(slide, cx - size*0.13, cy - size*0.13, size*0.26, size*0.26, glyph_color)
    return t
