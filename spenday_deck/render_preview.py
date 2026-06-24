"""Lightweight faithful renderer: generated .pptx -> per-slide PNGs + contact
sheet. Supports the shape vocabulary this deck uses (rect/roundRect/ellipse/
triangle/star, textboxes with mixed runs, pictures, connectors). For QA only."""
import io
import os
import glob
import math
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn

HERE = os.path.dirname(__file__)
PPTX = os.path.join(HERE, "SpenDay.pptx")
OUT = os.path.join(HERE, "preview")
os.makedirs(OUT, exist_ok=True)

SS = 2  # supersample
W, H = 1600 * SS, 900 * SS
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
PX_PER_IN = W / SLIDE_W_IN
EMU_PER_IN = 914400


def emu_px(v):
    return v / EMU_PER_IN * PX_PER_IN


# ---- fonts -----------------------------------------------------------------
FONT_FILES = {}
for d in [os.path.expanduser("~/Library/Fonts"), "/Library/Fonts",
          "/System/Library/Fonts", "/System/Library/Fonts/Supplemental"]:
    for f in glob.glob(os.path.join(d, "Montserrat*.ttf")):
        FONT_FILES[os.path.basename(f)] = f

NAME_TO_FILE = {
    "Montserrat Black": "Montserrat-Black.ttf",
    "Montserrat ExtraBold": "Montserrat-ExtraBold.ttf",
    "Montserrat SemiBold": "Montserrat-SemiBold.ttf",
    "Montserrat Medium": "Montserrat-Medium.ttf",
    "Montserrat": "Montserrat-Regular.ttf",
    "Montserrat Light": "Montserrat-Light.ttf",
}
_font_cache = {}


def get_font(name, size_pt):
    fn = NAME_TO_FILE.get(name, "Montserrat-Regular.ttf")
    path = FONT_FILES.get(fn)
    px = int(round(size_pt * PX_PER_IN / 72.0))
    key = (fn, px)
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(path, px) if path else \
            ImageFont.load_default()
    return _font_cache[key]


def rgb(color):
    return (color[0], color[1], color[2])


# ---- geometry helpers ------------------------------------------------------
def star_points(cx, cy, r):
    pts = []
    for i in range(10):
        ang = -math.pi / 2 + i * math.pi / 5
        rr = r if i % 2 == 0 else r * 0.42
        pts.append((cx + rr * math.cos(ang), cy + rr * math.sin(ang)))
    return pts


def draw_shape(draw, prst, box, fill, line, lw, adj):
    x0, y0, x1, y1 = box
    if prst in ("rect",):
        if fill:
            draw.rectangle(box, fill=fill)
        if line:
            draw.rectangle(box, outline=line, width=lw)
    elif prst == "roundRect":
        rad = adj * min(x1 - x0, y1 - y0)
        draw.rounded_rectangle(box, radius=rad, fill=fill, outline=line,
                               width=lw if line else 1)
    elif prst == "ellipse":
        draw.ellipse(box, fill=fill, outline=line, width=lw if line else 1)
    elif prst == "triangle":
        draw.polygon([(x0, y1), (x1, y1), ((x0 + x1) / 2, y0)], fill=fill)
    elif prst.startswith("star"):
        draw.polygon(star_points((x0 + x1) / 2, (y0 + y1) / 2,
                                  (x1 - x0) / 2), fill=fill)
    else:
        if fill:
            draw.rounded_rectangle(box, radius=6 * SS, fill=fill, outline=line)


# ---- text layout -----------------------------------------------------------
def get_runs(p):
    out = []
    for r in p.runs:
        col = (30, 27, 74)
        try:
            c = r.font.color.rgb
            col = (c[0], c[1], c[2])
        except Exception:
            pass
        out.append({
            "text": r.text,
            "font": r.font.name or "Montserrat",
            "size": r.font.size.pt if r.font.size else 14,
            "color": col,
        })
    return out


def layout_paragraph(runs, max_w):
    """-> list of lines; each line = (tokens, width, height)."""
    tokens = []  # (text, font, size, color, is_space)
    for run in runs:
        parts = run["text"].split("\n")
        for k, seg in enumerate(parts):
            if k > 0:
                tokens.append(("\n", run["font"], run["size"], run["color"]))
            i = 0
            for word in seg.split(" "):
                if i > 0:
                    tokens.append((" ", run["font"], run["size"], run["color"]))
                if word:
                    tokens.append((word, run["font"], run["size"], run["color"]))
                i += 1
    lines, cur, cur_w = [], [], 0.0

    def flush():
        nonlocal cur, cur_w
        h = max((t[2] for t in cur), default=12) * PX_PER_IN / 72.0 * 1.0
        lines.append((cur, cur_w, h))
        cur, cur_w = [], 0.0

    for tok in tokens:
        if tok[0] == "\n":
            flush()
            continue
        f = get_font(tok[1], tok[2])
        w = f.getlength(tok[0])
        if tok[0] == " ":
            cur.append(tok); cur_w += w; continue
        if cur_w + w > max_w and cur:
            # drop trailing space
            while cur and cur[-1][0] == " ":
                cur_w -= get_font(cur[-1][1], cur[-1][2]).getlength(" ")
                cur.pop()
            flush()
        cur.append(tok); cur_w += w
    flush()
    return lines


def draw_textframe(img, draw, shape):
    tf = shape.text_frame
    bx0 = emu_px(shape.left); by0 = emu_px(shape.top)
    bw = emu_px(shape.width); bh = emu_px(shape.height)
    anchor = tf.vertical_anchor
    paras = []
    total_h = 0.0
    for p in tf.paragraphs:
        runs = get_runs(p)
        if not runs:
            sz = 12
            lines = [([], 0, sz * PX_PER_IN / 72.0)]
        else:
            lines = layout_paragraph(runs, bw)
        ls = p.line_spacing or 1.0
        before = (p.space_before.pt if p.space_before else 0) * PX_PER_IN / 72.0
        after = (p.space_after.pt if p.space_after else 0) * PX_PER_IN / 72.0
        ph = before + sum(l[2] * ls for l in lines) + after
        align = p.alignment
        paras.append((lines, ls, before, after, align))
        total_h += ph
    astr = str(anchor).upper()
    if "MIDDLE" in astr:
        y = by0 + max(0, (bh - total_h) / 2)
    elif "BOTTOM" in astr:
        y = by0 + max(0, bh - total_h)
    else:
        y = by0
    for lines, ls, before, after, align in paras:
        y += before
        alstr = str(align).upper()
        for toks, lw_, lh in lines:
            line_h = lh * ls
            if "RIGHT" in alstr:
                x = bx0 + bw - lw_
            elif "CENTER" in alstr:
                x = bx0 + (bw - lw_) / 2
            else:
                x = bx0
            # baseline-ish: draw with ascent alignment per token
            for text, fname, size, color in toks:
                f = get_font(fname, size)
                asc, desc = f.getmetrics()
                draw.text((x, y + (lh - (asc + desc)) + 0), text, font=f,
                          fill=rgb(color))
                x += f.getlength(text)
            y += line_h
        y += after


# ---- main ------------------------------------------------------------------
def render_slide(slide, idx):
    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img, "RGBA")
    for sh in slide.shapes:
        tag = sh._element.tag
        if tag.endswith("}pic"):
            try:
                blob = sh.image.blob
                im = Image.open(io.BytesIO(blob)).convert("RGBA")
                box = (int(emu_px(sh.left)), int(emu_px(sh.top)))
                size = (int(emu_px(sh.width)), int(emu_px(sh.height)))
                im = im.resize(size)
                img.paste(im, box, im)
            except Exception as e:
                print("pic err", e)
            continue
        if tag.endswith("}cxnSp"):
            xfrm = sh._element.spPr.find(qn("a:xfrm"))
            off = xfrm.find(qn("a:off")); ext = xfrm.find(qn("a:ext"))
            ox, oy = int(off.get("x")), int(off.get("y"))
            ex, ey = int(ext.get("cx")), int(ext.get("cy"))
            flipH = xfrm.get("flipH") == "1"
            flipV = xfrm.get("flipV") == "1"
            x0, y0 = ox, oy
            x1, y1 = ox + ex, oy + ey
            if flipH:
                x0, x1 = x1, x0
            if flipV:
                y0, y1 = y1, y0
            col = (200, 200, 210)
            try:
                c = sh.line.color.rgb; col = (c[0], c[1], c[2])
            except Exception:
                pass
            wpx = max(1, int(emu_px(sh.line.width or 12700)))
            draw.line([(emu_px(x0), emu_px(y0)), (emu_px(x1), emu_px(y1))],
                      fill=col, width=wpx)
            continue
        if tag.endswith("}sp"):
            spPr = sh._element.spPr
            geom = spPr.find(qn("a:prstGeom"))
            prst = geom.get("prst") if geom is not None else "rect"
            fill = None
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    c = sh.fill.fore_color.rgb; fill = (c[0], c[1], c[2])
            except Exception:
                pass
            line = None; lw = 1
            try:
                if sh.line.color and sh.line.color.type is not None:
                    c = sh.line.color.rgb; line = (c[0], c[1], c[2])
                    lw = max(1, int(emu_px(sh.line.width or 12700)))
            except Exception:
                pass
            adj = 0.16
            try:
                if sh.adjustments and len(sh.adjustments):
                    adj = sh.adjustments[0]
            except Exception:
                pass
            box = (emu_px(sh.left), emu_px(sh.top),
                   emu_px(sh.left + sh.width), emu_px(sh.top + sh.height))
            if fill or line:
                draw_shape(draw, prst, box, fill, line, lw, adj)
            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_textframe(img, draw, sh)
            continue
    final = img.resize((1600, 900), Image.LANCZOS)
    final.save(os.path.join(OUT, f"slide_{idx:02d}.png"))
    return final


def main():
    prs = Presentation(PPTX)
    thumbs = []
    for i, slide in enumerate(prs.slides, 1):
        thumbs.append(render_slide(slide, i))
    # contact sheet 3 cols
    cols, pad = 3, 16
    tw, th = 520, 293
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * tw + (cols + 1) * pad,
                              rows * th + (rows + 1) * pad), (235, 235, 242))
    for i, t in enumerate(thumbs):
        r, c = divmod(i, cols)
        sheet.paste(t.resize((tw, th)),
                    (pad + c * (tw + pad), pad + r * (th + pad)))
    sheet.save(os.path.join(OUT, "_contact_sheet.png"))
    print("rendered", len(thumbs), "slides ->", OUT)


if __name__ == "__main__":
    main()
